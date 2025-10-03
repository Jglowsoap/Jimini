"""Generate an executive Jimini pitch deck as a PowerPoint (.pptx).

Usage (after installing requirements):

    python scripts/generate_ppt.py

Outputs: Jimini_Pitch.pptx in repo root (overwrite safe).

Slides included:
 1. Title
 2. Problem
 3. What Jimini Does
 4. Architecture (concept)
 5. Rules-as-Code
 6. Shadow → Enforce Rollout
 7. Evidence & Observability
 8. Security & Noise Reduction
 9. Value Summary
10. Pilot Success Metrics
11. Risks & Mitigations
12. Ask / Next Steps
13. One-Liner

Extend by editing the SLIDES list below.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt


SLIDES = [
    ("Jimini: AI Policy Firewall", "Inline enforcement • auditability • safe rollout"),
    (
        "Problem",
        "AI agents can leak secrets, PII, and regulated content. Manual review is slow and non-auditable.",
    ),
    (
        "What Jimini Does",
        "Evaluates every inbound/outbound message against YAML rules (regex, length, optional LLM, direction, endpoint). Returns block / flag / allow with deterministic precedence.",
    ),
    (
        "Architecture (Concept)",
        "Agent/App → Jimini → External. Components: FastAPI gateway, rules loader (hot reload), enforcement engine, hash-chained audit log, metrics + SARIF, optional webhook + OTEL spans, optional OpenAI check.",
    ),
    (
        "Rules-as-Code",
        "pattern + min_count • max_chars • llm_prompt (fail-safe) • applies_to • endpoints (exact/prefix/glob) • shadow_override • suppression of generic API-1.0 when specific secret matches.",
    ),
    (
        "Shadow → Enforce Rollout",
        "1) Shadow baseline 2) Enforce high-confidence secrets 3) Expand to regulated packs 4) SIEM dashboards (SARIF) 5) Ticketing integration (future).",
    ),
    (
        "Evidence & Observability",
        "/v1/metrics • /v1/audit/verify • /v1/audit/sarif • Webhooks on block • OTEL spans (optional) • Hash chain integrity.",
    ),
    (
        "Security & Noise Reduction",
        "Precedence block>flag>allow • Specific secret suppression • Fail-safe LLM path • Minimal surface (single POST) • Deterministic outputs.",
    ),
    (
        "Value Summary",
        "Risk reduction • Compliance evidence • Fast integration (one POST) • No model retraining • Quantifiable metrics for governance.",
    ),
    (
        "Pilot Success Metrics",
        "Violations/1K msgs • False positive rate <5% (enforced) • Time-to-detect (real-time) • Secret exposure reduction vs baseline.",
    ),
    (
        "Risks & Mitigations",
        "Over-blocking: shadow first • Rule sprawl: lint + packs • False positives: scoping/suppression • Missing patterns: metrics review.",
    ),
    (
        "Ask / Next Steps",
        "Approve 2-week shadow pilot (≤1 engineer day). Deliver baseline report + enforce list, then phased rollout.",
    ),
    (
        "One-Liner",
        "Jimini = seatbelt + black box for AI agent I/O: rules-as-code guardrails, safe rollout, immutable evidence.",
    ),
]


def create_deck(path: Path) -> None:
    prs = Presentation()

    for title, body in SLIDES:
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body_tf = slide.shapes.placeholders[1].text_frame
        # Split body into bullet lines by ' • ' or '. ' heuristically
        if " • " in body:
            bullets = [b.strip() for b in body.split(" • ") if b.strip()]
        else:
            # fallback: split sentences
            bullets = [
                seg.strip()
                for seg in body.replace(";", ". ").split(". ")
                if seg.strip()
            ]
        # First bullet as paragraph text
        if bullets:
            body_tf.text = bullets[0]
            for b in bullets[1:]:
                p = body_tf.add_paragraph()
                p.text = b
                p.level = 0
        else:
            body_tf.text = body
        # Light formatting
        for p in body_tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(20)

    prs.save(path)


if __name__ == "__main__":
    out = Path("Jimini_Pitch.pptx")
    create_deck(out)
    print(f"[Jimini] Wrote {out.resolve()}")
